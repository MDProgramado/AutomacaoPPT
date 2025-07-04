# -*- coding: utf-8 -*-
import os
import shutil
import logging
import cv2
import configparser
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm
from PIL import Image
import tkinter as tk
from tkinter import filedialog, scrolledtext, ttk
import threading
import queue
import json
import sys

# ===================================================================================
# 1. FUNÇÃO CENTRALIZADA E CORRETA PARA ENCONTRAR ARQUIVOS (A NOSSA "FERRAMENTA")
# ===================================================================================
def resource_path(relative_path):
    """ Retorna o caminho absoluto para o recurso, funcionando em dev e no PyInstaller. """
    try:
        # PyInstaller cria uma pasta temporária e armazena o caminho em sys._MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        # Se não estiver empacotado, o caminho é o do diretório do script
        base_path = os.path.abspath(".")
    
    return os.path.join(base_path, relative_path)

# ===================================================================================
# 2. FUNÇÃO ÚNICA E ROBUSTA PARA CARREGAR O ARQUIVO DE CONFIGURAÇÃO
# ===================================================================================
def carregar_configuracao():
    """ Carrega o config.ini usando o caminho correto e trata exceções. """
    config_path = resource_path('config.ini') # Usa a nossa ferramenta!
    try:
        config = configparser.ConfigParser()
        # É importante verificar a existência ANTES de tentar ler.
        if not os.path.exists(config_path):
            # Loga o erro e também lança uma exceção para parar a execução
            error_msg = f"Erro Crítico: O ficheiro de configuração 'config.ini' não foi encontrado no caminho: {config_path}"
            logging.critical(error_msg)
            raise FileNotFoundError(error_msg)
            
        config.read(config_path, encoding='utf-8')
        return config
    except Exception as e:
        logging.error(f"Erro ao carregar ou ler o arquivo de configuração: {e}")
        raise # Re-lança a exceção para que o programa pare se a config falhar.

# --- Configuração do Logging ---
# (Opcional, mas boa prática) Usar resource_path para o log também!
log_path = resource_path('log_automacao.txt')
logging.basicConfig(
    filename=log_path,
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    encoding='utf-8'
)

class AutomacaoPPT:
    # ===================================================================================
    # 3. CONSTRUTOR DA CLASSE SIMPLIFICADO
    # ===================================================================================
    def __init__(self, config):
        # A classe agora recebe a configuração já carregada, evitando redundância.
        self.config = config
        # Removemos a necessidade de carregar o config aqui dentro.

    # A função carregar_config foi REMOVIDA de dentro da classe para evitar duplicação.

    def verificar_desfocagem(self, caminho_imagem):
        try:
            limiar = self.config.getfloat('Configuracoes', 'limiar_desfocagem')
            imagem = cv2.imread(caminho_imagem)
            if imagem is None:
                return False, True
            
            cinza = cv2.cvtColor(imagem, cv2.COLOR_BGR2GRAY)
            variancia_laplaciano = cv2.Laplacian(cinza, cv2.CV_64F).var()
            
            if variancia_laplaciano < limiar:
                logging.warning(f"Imagem possivelmente desfocada: {os.path.basename(caminho_imagem)} (Variância: {variancia_laplaciano:.2f})")
                return True, False
            return False, False
        except Exception as e:
            logging.error(f"Não foi possível verificar a desfocagem para {os.path.basename(caminho_imagem)}: {e}")
            return False, True

    def processar_imagens(self, pasta_origem_sobrescrever=None, gui_queue=None):
        # O resto desta função pode continuar exatamente como estava.
        # Ela já usa self.config, que agora será carregado da maneira correta.
        # ... (todo o seu código de processar_imagens continua aqui, sem alterações)
        try:
            # Carrega as configurações das pastas
            pasta_origem = pasta_origem_sobrescrever if pasta_origem_sobrescrever else self.config['Pastas']['pasta_origem']
            pasta_destino = self.config['Pastas']['pasta_destino']
            pasta_processadas = self.config['Pastas']['pasta_processadas']
            
            # Use resource_path para o template também, por segurança!
            ficheiro_template_relativo = self.config['Pastas']['ficheiro_template']
            ficheiro_template = resource_path(ficheiro_template_relativo)
            
            # Garante que as pastas de destino existam
            os.makedirs(pasta_destino, exist_ok=True)
            os.makedirs(pasta_processadas, exist_ok=True)

            if not os.path.exists(pasta_origem):
                logging.error(f"A pasta de origem '{pasta_origem}' não existe.")
                if gui_queue: gui_queue.put(f"ERRO: Pasta de origem não encontrada: {pasta_origem}")
                return
            
            if not os.path.exists(ficheiro_template):
                logging.error(f"O ficheiro de template '{ficheiro_template}' não existe.")
                if gui_queue: gui_queue.put(f"ERRO: Ficheiro de template não encontrado: {ficheiro_template}")
                return

            # Carrega configurações de layout
            largura_cm = self.config.getfloat('Configuracoes', 'largura_cm')
            altura_cm = self.config.getfloat('Configuracoes', 'altura_cm')
            layout_por_slide = self.config.getint('Configuracoes', 'layout_por_slide')
            posicoes_str = self.config['Configuracoes']['posicoes']
            posicoes = json.loads(posicoes_str)

            # Abre a apresentação modelo
            prs = Presentation(ficheiro_template)

            # Adiciona o slide de cabeçalho
            slide_inicio = prs.slides.add_slide(prs.slide_layouts[0])
            titulo = slide_inicio.shapes.title
            titulo.text = "RELATÓRIO FOTOGRÁFICO"
            
            unidade = self.config['Pastas']['unidade']
            endereco = self.config['Pastas']['endereco']
            data = datetime.now().strftime("%d/%m/%Y")
            tipo_servico = "CORRETIVO"
            slide_inicio.shapes.placeholders[1].text = f"UNIDADE: {unidade}\nENDEREÇO: {endereco}\nCLASSIFICAÇÃO DO SERVIÇO: ({tipo_servico})\nDATA: {data}"

            imagens_encontradas = sorted([f for f in os.listdir(pasta_origem) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))])

            if not imagens_encontradas:
                logging.info("Nenhuma imagem encontrada na pasta de origem.")
                if gui_queue: gui_queue.put("Nenhuma imagem encontrada para processar.")
                return

            contador_imagens_no_slide = 0
            slide_atual = None
            total_imagens = len(imagens_encontradas)
            
            for i, nome_ficheiro in enumerate(imagens_encontradas):
                caminho_completo = os.path.join(pasta_origem, nome_ficheiro)
                
                if gui_queue:
                    progresso = int(((i + 1) / total_imagens) * 100)
                    gui_queue.put(f"PROGRESSO:{progresso}")
                    gui_queue.put(f"Processando {i+1}/{total_imagens}: {nome_ficheiro}")

                try:
                    with Image.open(caminho_completo) as img:
                        img.verify()
                    
                    desfocada, erro_leitura = self.verificar_desfocagem(caminho_completo)
                    if erro_leitura:
                        logging.error(f"Erro ao ler a imagem {nome_ficheiro} com OpenCV. Pulando.")
                        continue
                    
                    if contador_imagens_no_slide % layout_por_slide == 0:
                        template_slide_layout = prs.slide_layouts[5]
                        slide_atual = prs.slides.add_slide(template_slide_layout)
                        logging.info(f"Adicionando novo slide para as próximas {layout_por_slide} imagens.")
                    
                    posicao_atual = posicoes[contador_imagens_no_slide % layout_por_slide]
                    left = Cm(posicao_atual['left'])
                    top = Cm(posicao_atual['top'])
                    
                    slide_atual.shapes.add_picture(caminho_completo, left, top, width=Cm(largura_cm), height=Cm(altura_cm))
                    logging.info(f"Imagem '{nome_ficheiro}' adicionada ao slide.")
                    
                    contador_imagens_no_slide += 1
                    
                    shutil.move(caminho_completo, os.path.join(pasta_processadas, nome_ficheiro))

                except (IOError, SyntaxError) as e:
                    logging.error(f"Ficheiro '{nome_ficheiro}' está corrompido ou não é uma imagem válida: {e}")
                    if gui_queue: gui_queue.put(f"ERRO: Ficheiro corrompido: {nome_ficheiro}")
                except Exception as e:
                    logging.error(f"Erro inesperado ao processar '{nome_ficheiro}': {e}")
                    if gui_queue: gui_queue.put(f"ERRO inesperado com: {nome_ficheiro}")

            timestamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
            nome_final = f"Relatorio_Fotografico_{timestamp}.pptx"
            caminho_final = os.path.join(pasta_destino, nome_final)
            prs.save(caminho_final)
            
            logging.info(f"Relatório final salvo como '{caminho_final}'")
            if gui_queue:
                gui_queue.put("--------------------------------------------------")
                gui_queue.put(f"PROCESSO CONCLUÍDO COM SUCESSO!")
                gui_queue.put(f"Relatório salvo em: {caminho_final}")
                gui_queue.put("FINALIZADO")

        except Exception as e:
            logging.critical(f"Ocorreu um erro crítico no processo: {e}", exc_info=True)
            if gui_queue:
                gui_queue.put(f"ERRO CRÍTICO: {e}")
                gui_queue.put("FINALIZADO")

class App:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("Gerador de Relatório Fotográfico")
        self.root.geometry("700x500")

        # ===================================================================================
        # 4. CARREGAMOS A CONFIGURAÇÃO UMA ÚNICA VEZ AO INICIAR O APP
        # ===================================================================================
        try:
            # Carrega a configuração uma vez, no início de tudo.
            configuracao = carregar_configuracao()
            
            # Injeta a configuração na classe de automação.
            self.automacao = AutomacaoPPT(configuracao) 
            self.pasta_origem_var = tk.StringVar(value=self.automacao.config['Pastas']['pasta_origem'])
            self.setup_widgets()

        except FileNotFoundError as e:
            # Se o config.ini não for encontrado, o app não pode funcionar.
            # Mostra um erro claro para o usuário.
            self.setup_error_widgets(str(e))
        except Exception as e:
            self.setup_error_widgets(f"Erro ao inicializar: {e}")

    def setup_widgets(self):
        """Cria os widgets normais da aplicação."""
        frame = ttk.Frame(self.root, padding="10")
        frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        ttk.Label(frame, text="Pasta com as Imagens:").grid(row=0, column=0, sticky=tk.W, pady=2)
        self.entry_pasta = ttk.Entry(frame, textvariable=self.pasta_origem_var, width=60)
        self.entry_pasta.grid(row=1, column=0, sticky=(tk.W, tk.E), pady=2)
        self.btn_selecionar = ttk.Button(frame, text="Selecionar Pasta", command=self.selecionar_pasta)
        self.btn_selecionar.grid(row=1, column=1, sticky=tk.W, padx=5)
        self.btn_iniciar = ttk.Button(frame, text="Gerar Relatório", command=self.iniciar_processamento)
        self.btn_iniciar.grid(row=2, column=0, columnspan=2, pady=10)
        self.progresso = ttk.Progressbar(frame, orient="horizontal", length=500, mode="determinate")
        self.progresso.grid(row=3, column=0, columnspan=2, pady=5)
        self.log_area = scrolledtext.ScrolledText(frame, wrap=tk.WORD, width=80, height=20)
        self.log_area.grid(row=4, column=0, columnspan=2, pady=5)
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        frame.columnconfigure(0, weight=1)
        
    def setup_error_widgets(self, error_message):
        """Cria widgets para mostrar uma mensagem de erro fatal."""
        frame = ttk.Frame(self.root, padding="20")
        frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        label = ttk.Label(frame, text=error_message, foreground="red", wraplength=650)
        label.pack(pady=20)
        close_button = ttk.Button(frame, text="Fechar", command=self.root.destroy)
        close_button.pack(pady=10)
        
    # O resto da classe App continua como estava...
    # ... (selecionar_pasta, iniciar_processamento, verificar_queue, run) ...
    def selecionar_pasta(self):
        pasta = filedialog.askdirectory(title="Selecione a pasta com as imagens")
        if pasta:
            self.pasta_origem_var.set(pasta)

    def iniciar_processamento(self):
        self.btn_iniciar.config(state="disabled")
        self.log_area.delete('1.0', tk.END)
        self.progresso["value"] = 0
        
        self.queue = queue.Queue()
        self.thread = threading.Thread(
            target=self.automacao.processar_imagens,
            args=(self.pasta_origem_var.get(), self.queue)
        )
        self.thread.start()
        self.root.after(100, self.verificar_queue)

    def verificar_queue(self):
        try:
            while True:
                msg = self.queue.get_nowait()
                if msg.startswith("PROGRESSO:"):
                    valor = int(msg.split(":")[1])
                    self.progresso["value"] = valor
                elif msg == "FINALIZADO":
                    self.btn_iniciar.config(state="normal")
                    return
                else:
                    self.log_area.insert(tk.END, msg + '\n')
                    self.log_area.see(tk.END)
        except queue.Empty:
            pass
        
        if self.thread.is_alive():
            self.root.after(100, self.verificar_queue)
        else:
            self.btn_iniciar.config(state="normal")

    def run(self):
        self.root.mainloop()


if __name__ == "__main__":
    app = App()
    app.run()